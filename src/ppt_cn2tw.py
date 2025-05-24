#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 簡體中文轉繁體中文轉換工具

功能：
1. 將 PPTX 檔案中的簡體中文轉換為繁體中文
2. 保留原始格式和樣式
3. 支援批次處理目錄中的多個檔案
"""

import os
import sys
import argparse
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from opencc import OpenCC

class PPTConverter:
    def __init__(self, verbose: bool = False):
        """初始化轉換器"""
        self.cc = OpenCC('s2t')
        self.verbose = verbose
    
    def log(self, message: str) -> None:
        """輸出日誌訊息"""
        if self.verbose:
            print(f"[INFO] {message}")
    
    def convert_text(self, text: Optional[str]) -> Optional[str]:
        """將簡體中文轉換為繁體中文"""
        if not text or not text.strip():
            return text
        return self.cc.convert(text)
    
    def process_ppt(self, input_path: str, output_path: str) -> bool:
        """處理單個PPTX檔案"""
        try:
            self.log(f"正在處理: {input_path}")
            
            # 載入簡體中文PPT
            prs = Presentation(input_path)
            
            # 處理幻燈片中的文字
            for slide in prs.slides:
                # 處理文字方塊
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        self._process_text_frame(shape.text_frame)
                    
                    # 處理表格
                    if shape.has_table:
                        self._process_table(shape.table)
                
                # 處理幻燈片備註
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    self._process_text_frame(slide.notes_slide.notes_text_frame)
            
            # 確保輸出目錄存在
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # 儲存為繁體中文PPT
            prs.save(output_path)
            self.log(f"轉換成功: {output_path}")
            return True
            
        except Exception as e:
            print(f"[錯誤] 處理檔案 {input_path} 時發生錯誤: {str(e)}", file=sys.stderr)
            return False
    
    def _process_text_frame(self, text_frame) -> None:
        """處理文字框架"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.text = self.convert_text(run.text)
    
    def _process_table(self, table) -> None:
        """處理表格"""
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    self._process_text_frame(cell.text_frame)

def find_pptx_files(directory: str, recursive: bool = False) -> List[str]:
    """尋找目錄中的所有PPTX檔案"""
    path = Path(directory)
    if recursive:
        return [str(p) for p in path.rglob('*.pptx')]
    return [str(p) for p in path.glob('*.pptx')]

def ensure_output_path(input_path: str, output_path: str) -> str:
    """確保輸出路徑正確"""
    input_path = Path(input_path)
    output_path = Path(output_path)
    
    # 如果輸出路徑是目錄，則自動生成輸出檔案名稱
    if output_path.is_dir() or not output_path.suffix == '.pptx':
        output_path = output_path / f"{input_path.stem}_tw{input_path.suffix}"
    
    return str(output_path)

def parse_arguments():
    """解析命令列參數"""
    parser = argparse.ArgumentParser(description='將PPTX檔案從簡體中文轉換為繁體中文')
    
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument('-i', '--input', help='輸入檔案路徑')
    group.add_argument('-d', '--dir', help='輸入目錄路徑')
    
    parser.add_argument('-o', '--output', help='輸出檔案或目錄路徑', default='./output')
    parser.add_argument('-r', '--recursive', action='store_true', help='遞迴處理子目錄')
    parser.add_argument('-v', '--verbose', action='store_true', help='顯示詳細輸出')
    parser.add_argument('--version', action='version', version='%(prog)s 1.0.0')
    
    return parser.parse_args()

def main():
    # 解析命令列參數
    args = parse_arguments()
    
    # 初始化轉換器
    converter = PPTConverter(verbose=args.verbose)
    
    # 處理單個檔案
    if args.input:
        input_path = os.path.abspath(args.input)
        if not os.path.isfile(input_path):
            print(f"錯誤: 找不到輸入檔案: {input_path}", file=sys.stderr)
            sys.exit(1)
            
        output_path = os.path.abspath(args.output)
        output_path = ensure_output_path(input_path, output_path)
        
        success = converter.process_ppt(input_path, output_path)
        sys.exit(0 if success else 1)
    
    # 處理目錄
    if args.dir:
        input_dir = os.path.abspath(args.dir)
        if not os.path.isdir(input_dir):
            print(f"錯誤: 找不到輸入目錄: {input_dir}", file=sys.stderr)
            sys.exit(1)
        
        output_dir = os.path.abspath(args.output)
        os.makedirs(output_dir, exist_ok=True)
        
        # 尋找所有 PPTX 檔案
        pptx_files = find_pptx_files(input_dir, recursive=args.recursive)
        
        if not pptx_files:
            print(f"在 {input_dir} 中未找到任何 PPTX 檔案", file=sys.stderr)
            sys.exit(1)
        
        print(f"找到 {len(pptx_files)} 個 PPTX 檔案，開始轉換...")
        
        # 處理每個檔案
        success_count = 0
        for input_path in pptx_files:
            # 跳過已經轉換過的檔案（以 _tw.pptx 結尾）
            if '_tw.' in input_path.lower():
                print(f"跳過已轉換檔案: {input_path}")
                continue
                
            # 計算相對路徑以保持目錄結構
            rel_path = os.path.relpath(input_path, input_dir)
            output_path = os.path.join(output_dir, rel_path)
            output_path = ensure_output_path(input_path, output_path)
            
            # 確保輸出目錄存在
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            if converter.process_ppt(input_path, output_path):
                success_count += 1
        
        print(f"\n轉換完成！成功: {success_count}/{len(pptx_files)}")
        sys.exit(0 if success_count > 0 else 1)

if __name__ == "__main__":
    main()
