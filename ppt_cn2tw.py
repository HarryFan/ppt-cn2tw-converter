#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPT簡體中文轉繁體中文轉換工具

功能：
1. 遞迴處理指定目錄下的所有PPTX檔案
2. 將簡體中文轉換為繁體中文
3. 保留原始檔案，並在檔名後加上"_tw"作為新檔案
"""

import os
import sys
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from opencc import OpenCC

class PPTConverter:
    def __init__(self):
        # 初始化簡轉繁轉換器
        self.cc = OpenCC('s2t')
    
    def convert_text(self, text: Optional[str]) -> Optional[str]:
        """將簡體中文轉換為繁體中文"""
        if not text or not text.strip():
            return text
        return self.cc.convert(text)
    
    def process_ppt(self, input_path: str, output_path: str) -> bool:
        """處理單個PPTX檔案"""
        try:
            # 載入簡體中文PPT
            prs = Presentation(input_path)
            
            # 處理幻燈片中的文字
            for slide in prs.slides:
                # 處理文字方塊
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    # 處理段落
                    for paragraph in shape.text_frame.paragraphs:
                        # 直接處理每個 run 的文字，這樣可以保留原始格式
                        for run in paragraph.runs:
                            if run.text:
                                run.text = self.convert_text(run.text)
                    
                    # 處理表格
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    # 處理表格中的每個 run 以保留格式
                                    for run in paragraph.runs:
                                        if run.text:
                                            run.text = self.convert_text(run.text)
            
            # 處理幻燈片備註
            for slide in prs.slides:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text:
                        slide.notes_slide.notes_text_frame.text = self.convert_text(notes_text)
            
            # 儲存為繁體中文PPT
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"處理檔案 {input_path} 時發生錯誤: {str(e)}")
            return False

def find_pptx_files(directory: str) -> List[str]:
    """尋找目錄中的所有PPTX檔案"""
    return [str(p) for p in Path(directory).rglob('*.pptx')]

def main():
    import argparse
    
    # 設定參數解析器
    parser = argparse.ArgumentParser(description='PPT 簡體中文轉繁體中文轉換工具')
    parser.add_argument('directory', help='要處理的目錄路徑')
    parser.add_argument('-o', '--output', help='輸出目錄 (預設: 與來源目錄相同)')
    
    # 解析參數
    args = parser.parse_args()
    
    directory = args.directory
    output_dir = args.output if args.output else directory
    
    # 檢查目錄是否存在
    if not os.path.isdir(directory):
        print(f"錯誤: 目錄不存在: {directory}")
        sys.exit(1)
        
    # 如果指定了輸出目錄，確保它存在
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    
    # 建立轉換器
    converter = PPTConverter()
    
    # 尋找所有PPTX檔案
    pptx_files = find_pptx_files(directory)
    
    if not pptx_files:
        print(f"在 {directory} 中未找到任何PPTX檔案")
        return
    
    print(f"找到 {len(pptx_files)} 個PPTX檔案，開始轉換...")
    
    # 處理每個檔案
    success_count = 0
    for input_path in pptx_files:
        # 跳過已經轉換過的檔案（以 _tw.pptx 結尾）
        if input_path.endswith('_tw.pptx'):
            continue
            
        # 設定輸出檔案路徑
        path_obj = Path(input_path)
        
        # 處理檔名，移除可能已存在的 _tw_v2 後綴
        stem = path_obj.stem
        if stem.endswith('_tw_v2'):
            base_stem = stem[:-6]  # 移除 _tw_v2
        elif stem.endswith('_tw'):
            base_stem = stem[:-3]   # 移除 _tw
        else:
            base_stem = stem
            
        # 設定新的檔名
        new_stem = f"{base_stem}_tw"
        
        # 如果指定了輸出目錄，則使用該目錄
        if output_dir and output_dir != directory:
            # 保持原始子目錄結構
            rel_path = os.path.relpath(os.path.dirname(input_path), directory)
            target_dir = os.path.join(output_dir, rel_path)
            os.makedirs(target_dir, exist_ok=True)
            output_path = os.path.normpath(os.path.join(target_dir, f"{new_stem}{path_obj.suffix}"))
        else:
            # 使用原始目錄
            output_path = str(path_obj.with_stem(new_stem))
        
        print(f"正在處理: {input_path} -> {output_path}")
        
        # 執行轉換
        if converter.process_ppt(input_path, output_path):
            success_count += 1
            print(f"轉換成功: {output_path}")
        else:
            print(f"轉換失敗: {input_path}")
    
    print(f"\n轉換完成！成功: {success_count}/{len(pptx_files)}")

if __name__ == "__main__":
    main()
